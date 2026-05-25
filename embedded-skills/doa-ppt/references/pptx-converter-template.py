# -*- coding: utf-8 -*-
"""
python-pptx 单页演示文稿转换模板

本文件提供一套通用的工具函数和使用示例，用于将 HTML 预览内容转为 .pptx。
实际使用时，根据 HTML 的具体布局和内容定制颜色常量、布局参数和元素组装逻辑。

依赖: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ════════════════════════════════════════════
# 颜色常量模板（根据 HTML CSS 变量定义）
# ════════════════════════════════════════════
#
# 从 HTML 的 :root CSS 变量提取，转为 RGBColor：
#   --primary: #7C5CFC  →  PRIMARY = RGBColor(0x7C, 0x5C, 0xFC)
#   --primary-light: #A78BFA  →  PRIMARY_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
#
# 常用色彩角色：
#   PRIMARY      - 主题色（标题、按钮、强调）
#   PRIMARY_LIGHT - 浅色变体（卡片边框、图标底色）
#   ACCENT       - 点缀色
#   BG_SLIDE     - 幻灯片背景
#   BG_CARD      - 卡片背景（通常白色）
#   TEXT_PRIMARY  - 主文字色
#   TEXT_SECONDARY - 次要文字色
#   TEXT_MUTED    - 辅助说明色
#   BANNER_COLOR  - 横幅/按钮实色（渐变的中间值）
#   WHITE         - 白色

# 示例（紫色系浅色主题）：
PRIMARY = RGBColor(0x7C, 0x5C, 0xFC)
PRIMARY_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT = RGBColor(0xC0, 0x84, 0xFC)
BG_SLIDE = RGBColor(0xF5, 0xF0, 0xFF)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY = RGBColor(0x1E, 0x1B, 0x3A)
TEXT_SECONDARY = RGBColor(0x6B, 0x64, 0x94)
TEXT_MUTED = RGBColor(0x9B, 0x95, 0xB8)
DIVIDER_COLOR = RGBColor(0xF0, 0xEC, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BANNER_COLOR = RGBColor(0x8B, 0x5C, 0xF6)

# 幻灯片尺寸（16:9 宽屏）
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ════════════════════════════════════════════
# 工具函数
# ════════════════════════════════════════════

def set_slide_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, border_width=Pt(1), radius=Inches(0.15)):
    """
    添加圆角矩形

    Args:
        radius: 圆角半径，Inches/Pt 值。较大值 = 更圆的角。
                Inches(0.08) = 微圆角, Inches(0.15) = 常规, Inches(0.25) = 大圆角
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 设置圆角半径
    sp = shape._element
    prstGeom = None
    for child in sp.iter():
        if child.tag.endswith('}prstGeom'):
            prstGeom = child
            break
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = prstGeom.makeelement(qn('a:avLst'), {})
            prstGeom.append(avLst)
        for gd in list(avLst):
            avLst.remove(gd)
        r_val = int(radius / min(width, height) * 50000)
        r_val = min(r_val, 50000)
        gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': f'val {r_val}'})
        avLst.append(gd)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei', anchor='t', line_spacing=None):
    """
    添加文本框

    Args:
        anchor: 垂直对齐 - 't'=顶部, 'ctr'=居中, 'b'=底部
        line_spacing: 行间距（Pt 值），None 使用默认
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    txBody = txBox._element.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if txBody is not None:
        txBody.set('anchor', anchor)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=12,
                       color=TEXT_PRIMARY, font_name='Microsoft YaHei',
                       alignment=PP_ALIGN.LEFT, line_spacing=18, bold=False):
    """添加多行文本（每个字符串元素一个段落）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_rich_paragraph(slide, left, top, width, height, runs,
                       alignment=PP_ALIGN.LEFT, line_spacing=None):
    """
    添加富文本段落（单段多 run，支持混合粗体/颜色）

    Args:
        runs: list of dict, 每个 dict 支持：
            text (str), size (int, pt), color (RGBColor),
            bold (bool), font (str)

    示例:
        add_rich_paragraph(slide, x, y, w, h, [
            {'text': '报表名称: ', 'size': 10, 'color': TEXT_PRIMARY, 'bold': True},
            {'text': 'Sales Report', 'size': 10, 'color': PRIMARY},
        ])
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, run_def in enumerate(runs):
        if i == 0:
            r = p.runs[0] if p.runs else p.add_run()
        else:
            r = p.add_run()
        r.text = run_def.get('text', '')
        r.font.size = Pt(run_def.get('size', 12))
        r.font.color.rgb = run_def.get('color', TEXT_PRIMARY)
        r.font.bold = run_def.get('bold', False)
        r.font.name = run_def.get('font', 'Microsoft YaHei')
    return txBox


def add_circle(slide, left, top, size, fill_color):
    """添加圆形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_thin_rect(slide, left, top, width, height, fill_color):
    """添加极窄矩形（用于分隔线、色条装饰）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════
# 使用示例：组装一页完整幻灯片
# ════════════════════════════════════════════
#
# 以下为典型的单页 PPT 组装流程，包含：
# - Header（Logo + 标题 + 分隔线）
# - Left Section（功能描述 + 特性卡片）
# - Right Section（对话/交互预览）
# - Footer（指标横幅）
#
# 实际使用时根据 HTML 内容定制以下代码。

def build_example_slide():
    """示例：构建一页完整的产品介绍幻灯片"""

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, BG_SLIDE)

    MARGIN_X = Inches(0.7)

    # ── Header ──
    # Logo 圆角色块
    add_rounded_rect(slide, MARGIN_X, Inches(0.4), Inches(0.5), Inches(0.5),
                     BANNER_COLOR, radius=Inches(0.1))
    add_text_box(slide, MARGIN_X, Inches(0.4), Inches(0.5), Inches(0.5),
                 '⚡', font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER, anchor='ctr')
    # 标题
    add_text_box(slide, MARGIN_X + Inches(0.65), Inches(0.44), Inches(6), Inches(0.5),
                 '产品名称 · 功能介绍', font_size=26, color=PRIMARY, bold=True)
    # 分隔线
    add_thin_rect(slide, MARGIN_X, Inches(1.05), Inches(11.9), Pt(1.5),
                  RGBColor(0xE0, 0xD8, 0xF8))

    # ── Left Section: 特性卡片 ──
    LEFT_W = Inches(5.8)
    features = [
        ('💬', '特性一', '特性一的描述文字', RGBColor(0xED, 0xE9, 0xFE)),
        ('🎯', '特性二', '特性二的描述文字', RGBColor(0xE0, 0xE7, 0xFF)),
    ]
    card_top = Inches(2.0)
    for i, (icon, title, desc, icon_bg) in enumerate(features):
        cy = card_top + i * Inches(0.85)
        add_rounded_rect(slide, MARGIN_X, cy, LEFT_W, Inches(0.7), BG_CARD, radius=Inches(0.12))
        # 左侧色条
        add_thin_rect(slide, MARGIN_X + Pt(2), cy + Pt(8), Pt(3), Inches(0.7) - Pt(16),
                      PRIMARY_LIGHT)
        # 图标
        add_rounded_rect(slide, MARGIN_X + Inches(0.25), cy + Inches(0.14),
                         Inches(0.42), Inches(0.42), icon_bg, radius=Inches(0.08))
        add_text_box(slide, MARGIN_X + Inches(0.25), cy + Inches(0.14),
                     Inches(0.42), Inches(0.42), icon, font_size=16,
                     alignment=PP_ALIGN.CENTER, anchor='ctr')
        # 标题 + 描述
        add_text_box(slide, MARGIN_X + Inches(0.8), cy + Pt(8),
                     LEFT_W - Inches(1.0), Inches(0.25),
                     title, font_size=13, color=TEXT_PRIMARY, bold=True)
        add_text_box(slide, MARGIN_X + Inches(0.8), cy + Pt(26),
                     LEFT_W - Inches(1.0), Inches(0.3),
                     desc, font_size=10, color=TEXT_MUTED)

    # ── Footer: 指标横幅 ──
    footer_x = MARGIN_X
    footer_w = SLIDE_W - MARGIN_X * 2
    add_rounded_rect(slide, footer_x, Inches(6.45), footer_w, Inches(0.75),
                     BANNER_COLOR, radius=Inches(0.15))
    metrics = [('99%', '可用率'), ('500+', '用户数'), ('24/7', '在线')]
    metric_w = footer_w / len(metrics)
    for i, (val, lbl) in enumerate(metrics):
        mx = footer_x + i * metric_w
        add_text_box(slide, mx, Inches(6.49), metric_w, Inches(0.3),
                     val, font_size=22, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, mx, Inches(6.75), metric_w, Inches(0.25),
                     lbl, font_size=10, color=RGBColor(0xDD, 0xCC, 0xFF),
                     alignment=PP_ALIGN.CENTER)

    return prs


if __name__ == '__main__':
    prs = build_example_slide()
    prs.save('example_output.pptx')
    print('示例 PPT 已生成: example_output.pptx')
