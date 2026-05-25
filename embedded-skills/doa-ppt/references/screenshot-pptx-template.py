# -*- coding: utf-8 -*-
"""
截图高保真 PPT 生成模板

将可翻页 HTML 幻灯片逐页截图为高清 PNG，再插入 python-pptx 生成 .pptx。
实际使用时只需修改底部的 CONFIG 参数即可。

依赖:
    pip install python-pptx playwright
    playwright install chromium
"""

import os
import sys
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


# ════════════════════════════════════════════
# 常量
# ════════════════════════════════════════════

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 视口尺寸（与 HTML 中 .slides-viewport 一致）
VIEWPORT_W = 1280
VIEWPORT_H = 720

# 高清倍率（2 = Retina 级别，截图实际为 2560×1440）
DEVICE_SCALE = 2

# 翻页后等待过渡动画的时间（毫秒）
TRANSITION_WAIT = 600


# ════════════════════════════════════════════
# 核心函数
# ════════════════════════════════════════════

def capture_slides(html_path, output_dir, total_slides):
    """
    使用 Playwright 打开 HTML 幻灯片并逐页截图。

    Args:
        html_path:    slides.html 文件路径（相对或绝对）
        output_dir:   截图输出目录
        total_slides: 幻灯片总页数

    Returns:
        截图文件路径列表
    """
    abs_html = Path(html_path).resolve()
    if not abs_html.exists():
        print(f'错误: 找不到 HTML 文件 {abs_html}')
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)
    screenshots = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={'width': VIEWPORT_W, 'height': VIEWPORT_H},
            device_scale_factor=DEVICE_SCALE,
        )
        page.goto(f'file://{abs_html}')
        page.wait_for_load_state('networkidle')

        # 隐藏导航控件（箭头、页码圆点、键盘提示），避免遮挡内容
        page.evaluate('''() => {
            document.querySelectorAll('.nav-btn, .page-indicator, .keyboard-hint')
                .forEach(el => el.style.display = 'none');
        }''')

        for i in range(total_slides):
            # 通过 JS 全局函数 goTo() 切换到第 i 页
            page.evaluate(f'goTo({i})')
            page.wait_for_timeout(TRANSITION_WAIT)

            # 截图 .slides-viewport 容器
            viewport = page.query_selector('.slides-viewport')
            if viewport is None:
                print(f'警告: 未找到 .slides-viewport，跳过第 {i+1} 页')
                continue

            screenshot_path = os.path.join(output_dir, f'slide_{i+1:02d}.png')
            viewport.screenshot(path=screenshot_path)
            screenshots.append(screenshot_path)
            print(f'  ✓ 截图第 {i+1}/{total_slides} 页 → {screenshot_path}')

        browser.close()

    return screenshots


def build_pptx(screenshots, output_path):
    """
    将截图列表插入 PPT，每张图占满整页。

    Args:
        screenshots: 截图文件路径列表（按顺序）
        output_path: 输出 .pptx 路径
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for img_path in screenshots:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(
            img_path,
            left=0, top=0,
            width=SLIDE_W, height=SLIDE_H,
        )

    prs.save(output_path)
    print(f'\n截图高保真 PPT 已生成: {output_path}')


def auto_detect_total(html_path):
    """
    自动检测 HTML 中的幻灯片总数（通过查询 .slide 元素个数）。

    Args:
        html_path: slides.html 文件路径

    Returns:
        幻灯片总页数
    """
    abs_html = Path(html_path).resolve()
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={'width': VIEWPORT_W, 'height': VIEWPORT_H},
        )
        page.goto(f'file://{abs_html}')
        page.wait_for_load_state('networkidle')
        total = page.evaluate('document.querySelectorAll(".slide").length')
        browser.close()
    return total


# ════════════════════════════════════════════
# 使用入口 — 只需修改以下 CONFIG
# ════════════════════════════════════════════

if __name__ == '__main__':
    # ── CONFIG（每次按实际情况修改） ──
    HTML_FILE = 'slides.html'          # HTML 幻灯片文件路径
    OUTPUT_PPTX = 'output.pptx'        # 输出 PPT 路径
    SCREENSHOT_DIR = '_screenshots'    # 截图临时目录
    TOTAL_SLIDES = None                # None = 自动检测；或填写固定数字

    # ── 执行 ──
    if TOTAL_SLIDES is None:
        print('正在检测幻灯片总数...')
        TOTAL_SLIDES = auto_detect_total(HTML_FILE)
        print(f'检测到 {TOTAL_SLIDES} 页幻灯片')

    print(f'开始截图（{DEVICE_SCALE}x 高清）...')
    imgs = capture_slides(HTML_FILE, SCREENSHOT_DIR, TOTAL_SLIDES)

    if imgs:
        build_pptx(imgs, OUTPUT_PPTX)
    else:
        print('未生成任何截图，请检查 HTML 文件。')
